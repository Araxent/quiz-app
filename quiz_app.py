import streamlit as st
from pptx import Presentation
from io import BytesIO

st.set_page_config(page_title="Quiz Interactif", layout="centered")

# -------- CONFIGURATION --------
ADMIN_PASSWORD = "admin123"

# -------- SESSION INIT --------
if "role" not in st.session_state:
    st.session_state.role = None
if "pseudo" not in st.session_state:
    st.session_state.pseudo = ""
if "started" not in st.session_state:
    st.session_state.started = False
if "index" not in st.session_state:
    st.session_state.index = 0
if "score" not in st.session_state:
    st.session_state.score = 0
if "questions" not in st.session_state:
    st.session_state.questions = []
if "reponses_joueurs" not in st.session_state:
    st.session_state.reponses_joueurs = {}

# -------- CHARGEMENT QUESTIONS --------
@st.cache_data
def charger_questions(uploaded_file):
    prs = Presentation(uploaded_file)
    questions = []
    for slide in prs.slides:
        question = slide.shapes.title.text if slide.shapes.title else ""
        reponse = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                reponse = shape.text
                break
        if question and reponse:
            questions.append((question.strip(), reponse.strip()))
    return questions


# -------- PAGE SCORE FINAL --------
def page_score_finale():
    st.title("🏁 Résultats finaux du quiz")
    st.markdown("## 🏆 Classement des joueurs")

    scores = {
        pseudo: st.session_state.get(f"score_{pseudo}", 0)
        for pseudo in st.session_state.reponses_joueurs.keys()
    }
    scores_tries = dict(sorted(scores.items(), key=lambda item: item[1], reverse=True))

    for i, (pseudo, score) in enumerate(scores_tries.items(), 1):
        st.markdown(f"**{i}. {pseudo}** — {score} points")

    st.download_button("📥 Télécharger les scores en CSV",
                       data="Pseudo,Score\n" + "\n".join([f"{k},{v}" for k, v in scores_tries.items()]),
                       file_name="resultats_quiz.csv",
                       mime="text/csv")

# -------- INTERFACE ADMIN --------
def admin_interface():
    st.title("👩‍💻 Interface Admin")
    uploaded_file = st.file_uploader("📂 Chargez le fichier PowerPoint", type=["pptx"])
    if uploaded_file:
        st.session_state.questions = charger_questions(uploaded_file)
        st.success("✅ Questions chargées")

    if st.session_state.questions:
            if st.button("🏁 Terminer le quiz et voir les scores finaux"):
                st.session_state.index = len(st.session_state.questions)
                st.session_state.show_results = True

        if not st.session_state.started:
            if st.button("🚀 Lancer le quiz pour tous"):
                st.session_state.started = True
                st.success("Quiz lancé !")
        elif not st.session_state.get("show_results", False):
            question, reponse = st.session_state.questions[st.session_state.index]
            st.markdown(f"### Question {st.session_state.index + 1} : {question}")
            st.markdown(f"**Réponse attendue :** {reponse}")

            if st.button("⏭️ Question suivante"):
                if st.session_state.index + 1 < len(st.session_state.questions):
                    st.session_state.index += 1
                else:
                    st.info("Fin du quiz")

            st.markdown("---")
            st.markdown("### 🧠 Réponses des joueurs")
            for pseudo, data in st.session_state.reponses_joueurs.items():
                rep = data.get(st.session_state.index, "⏳ Pas encore répondu")
                st.write(f"**{pseudo}** : {rep}")

            st.markdown("---")
            st.markdown("### 🎯 Scores")
            for pseudo, data in st.session_state.reponses_joueurs.items():
                if f"score_{pseudo}" not in st.session_state:
                    st.session_state[f"score_{pseudo}"] = 0
                col1, col2, col3 = st.columns([1, 2, 1])
                with col1:
                    if st.button(f"+1 {pseudo}"):
                        st.session_state[f"score_{pseudo}"] += 1
                with col3:
                    if st.button(f"-1 {pseudo}"):
                        st.session_state[f"score_{pseudo}"] -= 1
                with col2:
                    st.markdown(f"<h4 style='text-align: center;'>{pseudo} : {st.session_state[f'score_{pseudo}']}</h4>", unsafe_allow_html=True)

# -------- INTERFACE JOUEUR --------
def joueur_interface():
    st.title("🎮 Quiz")
    if not st.session_state.started:
        st.info("⏳ En attente que l'admin lance le quiz…")
        st.stop()

    if st.session_state.index < len(st.session_state.questions):
        question, _ = st.session_state.questions[st.session_state.index]
        st.markdown(f"### Question {st.session_state.index + 1} : {question}")

        user_answer = st.text_input("✏️ Votre réponse", key=f"rep_{st.session_state.index}")
        if user_answer:
            if st.session_state.pseudo not in st.session_state.reponses_joueurs:
                st.session_state.reponses_joueurs[st.session_state.pseudo] = {}
            st.session_state.reponses_joueurs[st.session_state.pseudo][st.session_state.index] = user_answer

    st.markdown("---")
    st.markdown("### 🧮 Votre score")
    score = st.session_state.get(f"score_{st.session_state.pseudo}", 0)
    st.markdown(f"<h2 style='text-align: center;'>{score}</h2>", unsafe_allow_html=True)

# -------- INTERFACE DE CONNEXION --------
def login_interface():
    st.title("🔐 Connexion au Quiz")

    tab1, tab2 = st.tabs(["👤 Joueur", "🛠️ Admin"])

    with tab1:
        pseudo = st.text_input("Votre pseudo", key="pseudo_input")
        if st.button("Entrer dans la salle d'attente"):
            if pseudo:
                st.session_state.pseudo = pseudo
                st.session_state.role = "joueur"
                st.experimental_rerun()

    with tab2:
        mdp = st.text_input("Mot de passe admin", type="password")
        if st.button("Connexion admin"):
            if mdp == ADMIN_PASSWORD:
                st.session_state.role = "admin"
                st.experimental_rerun()
            else:
                st.error("Mot de passe incorrect.")


# -------- REDIRECTION SCORE FINAL --------
if st.session_state.role == "admin" and st.session_state.get("show_results", False):
    page_score_finale()
    st.stop()
    
# -------- ROUTING --------

if st.session_state.role == "admin":
    admin_interface()
elif st.session_state.role == "joueur":
    joueur_interface()
else:
    login_interface()
